"""Isolated product acceptance for provider-free current-document inspection."""

from __future__ import annotations

import json
from pathlib import Path
from types import SimpleNamespace
import zipfile

import fitz  # type: ignore[import-untyped]
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import pytest
from typer.testing import CliRunner

from local_steward.cli import app
from local_steward.config import load_config
from local_steward.document_observation import (
    DocumentInspectionRequest,
    inspect_document,
)
from local_steward.errors import (
    DocumentInspectionConfirmationError,
    DocumentInspectionInputError,
    DocumentInspectionScopeError,
    DocumentInspectionSourceChangedError,
)
from local_steward.file_agent.runtime import (
    ProjectOwnedBoundedDocumentIngress,
    ScopeBinding,
    ScopeBindings,
    StructuredDocumentParserAdapter,
)
from local_steward.models import StewardConfig


def _config(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    scopes: list[tuple[str, str, Path, bool]] | None = None,
) -> tuple[Path, StewardConfig]:
    monkeypatch.setattr("local_steward.scopes.SYSTEM_PROTECTED_PATHS", ())
    project = tmp_path / "project"
    config_path = project / "config" / "steward.toml"
    config_path.parent.mkdir(parents=True)
    values = scopes or [("managed", "managed_root", tmp_path / "documents", True)]
    scope_text = "\n".join(
        (
            "[[scopes]]\n"
            f'scope_id = "{scope_id}"\n'
            f'role = "{role}"\n'
            f'path = "{root}"\n'
            f"enabled = {'true' if enabled else 'false'}\n"
            "follow_directory_symlinks = false\n"
            "allow_cross_mount = false"
        )
        for scope_id, role, root, enabled in values
    )
    config_path.write_text(
        "\n".join(
            (
                "schema_version = 1",
                'project_name = "Document product test"',
                "[paths]",
                'data_dir = "data"',
                'cache_dir = "data/cache"',
                'evidence_dir = "data/evidence"',
                'quarantine_dir = "data/quarantine"',
                scope_text,
            )
        ),
        encoding="utf-8",
    )
    for _scope_id, _role, root, _enabled in values:
        root.mkdir(parents=True, exist_ok=True)
    return config_path, load_config(config_path)


def _write_pdf(path: Path, pages: int = 1) -> None:
    document = fitz.open()
    for index in range(pages):
        page = document.new_page()
        page.insert_text((72, 72), f"searchable PDF page {index + 1}: 文档事实")
    document.save(path)
    document.close()


def _write_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Facts"
    sheet.append(["Label", "Value"])
    sheet.append(["answer", 42])
    workbook.save(path)
    workbook.close()


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text_frame.text = "PPTX fact marker"
    presentation.save(path)


def _write_docx(path: Path) -> None:
    content_type = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" '
            'ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            f'<Override PartName="/word/document.xml" ContentType="{content_type}"/>'
            '<Override PartName="/word/styles.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>'
            "</Types>",
        )
        archive.writestr(
            "_rels/.rels",
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
            'Target="word/document.xml"/></Relationships>',
        )
        archive.writestr(
            "word/styles.xml",
            '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>',
        )
        archive.writestr(
            "word/document.xml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body><w:p><w:r><w:t>DOCX fact marker</w:t></w:r></w:p>"
            "<w:sectPr/></w:body></w:document>",
        )


def _write_epub(path: Path) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        archive.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?>'
            '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
            '<rootfiles><rootfile full-path="OEBPS/content.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        archive.writestr(
            "OEBPS/content.opf",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<package version="3.0" unique-identifier="book-id" '
            'xmlns="http://www.idpf.org/2007/opf">'
            '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            '<dc:identifier id="book-id">steward-test</dc:identifier>'
            '<dc:title>STEWARD EPUB</dc:title><dc:language>en</dc:language></metadata>'
            '<manifest><item id="chapter" href="chapter.xhtml" '
            'media-type="application/xhtml+xml"/></manifest>'
            '<spine><itemref idref="chapter"/></spine></package>',
        )
        archive.writestr(
            "OEBPS/chapter.xhtml",
            '<html xmlns="http://www.w3.org/1999/xhtml"><head><title>Chapter One</title></head>'
            '<body><h1>Chapter One</h1><p>EPUB fact marker</p></body></html>',
        )


def test_public_service_runs_all_five_real_workers(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    _path, config = _config(tmp_path, monkeypatch)
    root = tmp_path / "documents"
    writers = {
        "sample.pdf": _write_pdf,
        "sample.docx": _write_docx,
        "sample.xlsx": _write_xlsx,
        "sample.pptx": _write_pptx,
        "sample.epub": _write_epub,
    }
    expected = {
        "sample.pdf": ("PDF", "PyMuPDF4LLM"),
        "sample.docx": ("DOCX", "MarkItDown"),
        "sample.xlsx": ("XLSX", "openpyxl"),
        "sample.pptx": ("PPTX", "python-pptx"),
        "sample.epub": ("EPUB", "Docling"),
    }

    for name, writer in writers.items():
        writer(root / name)
        page = inspect_document(config, DocumentInspectionRequest("managed", name, True))
        assert page.status == "COMPLETE", name
        assert (page.source_format, page.backend_name) == expected[name]
        assert page.items
        assert page.source_sha256 is not None
        assert page.document_observation_digest is not None

    for name, query in {
        "sample.pdf": "SEARCHABLE pdf",
        "sample.docx": "DOCX fact marker",
        "sample.xlsx": "answer",
        "sample.pptx": "PPTX fact marker",
        "sample.epub": "EPUB fact marker",
    }.items():
        located = inspect_document(
            config,
            DocumentInspectionRequest("managed", name, True, content_query=query),
        )
        assert located.content_search is not None
        assert located.content_search.status == "COMPLETE"
        assert located.content_search.matched_item_count >= 1
        assert located.content_search.matches[0].location


def test_confirmation_scope_path_and_exclusion_are_typed(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    excluded = tmp_path / "documents" / "private"
    _path, config = _config(
        tmp_path,
        monkeypatch,
        [
            ("managed", "managed_root", tmp_path / "documents", True),
            ("reference", "reference_root", tmp_path / "reference", True),
            ("disabled", "managed_root", tmp_path / "disabled", False),
            ("excluded", "excluded_root", excluded, True),
        ],
    )
    _write_pdf(tmp_path / "documents" / "sample.pdf")
    _write_pdf(tmp_path / "reference" / "sample.pdf")
    _write_pdf(excluded / "secret.pdf")

    with pytest.raises(DocumentInspectionConfirmationError):
        inspect_document(config, DocumentInspectionRequest("managed", "sample.pdf", False))
    with pytest.raises(DocumentInspectionScopeError):
        inspect_document(config, DocumentInspectionRequest("unknown", "sample.pdf", True))
    with pytest.raises(DocumentInspectionScopeError):
        inspect_document(config, DocumentInspectionRequest("disabled", "sample.pdf", True))
    with pytest.raises(DocumentInspectionScopeError):
        inspect_document(config, DocumentInspectionRequest("excluded", "secret.pdf", True))
    with pytest.raises(DocumentInspectionScopeError):
        inspect_document(config, DocumentInspectionRequest("managed", "private/secret.pdf", True))
    with pytest.raises(DocumentInspectionInputError):
        inspect_document(config, DocumentInspectionRequest("managed", "../sample.pdf", True))

    reference = inspect_document(config, DocumentInspectionRequest("reference", "sample.pdf", True))
    assert reference.status == "COMPLETE"


def test_pagination_digest_and_source_pin_are_coherent(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    _path, config = _config(tmp_path, monkeypatch)
    source = tmp_path / "documents" / "sample.pdf"
    _write_pdf(source, pages=3)

    first = inspect_document(
        config, DocumentInspectionRequest("managed", "sample.pdf", True, limit=1)
    )
    assert first.returned_count == 1 and first.has_more and first.next_offset == 1
    assert first.source_sha256 is not None
    second = inspect_document(
        config,
        DocumentInspectionRequest(
            "managed",
            "sample.pdf",
            True,
            limit=1,
            offset=1,
            expected_source_sha256=first.source_sha256,
        ),
    )
    assert second.returned_count == 1
    assert second.document_observation_digest == first.document_observation_digest

    with pytest.raises(DocumentInspectionInputError):
        inspect_document(
            config, DocumentInspectionRequest("managed", "sample.pdf", True, limit=1, offset=1)
        )
    _write_pdf(source, pages=4)
    with pytest.raises(DocumentInspectionSourceChangedError):
        inspect_document(
            config,
            DocumentInspectionRequest(
                "managed",
                "sample.pdf",
                True,
                limit=1,
                offset=1,
                expected_source_sha256=first.source_sha256,
            ),
        )


def test_controlled_content_location_is_bounded_and_source_pinned(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    _path, config = _config(tmp_path, monkeypatch)
    source = tmp_path / "documents" / "sample.pdf"
    _write_pdf(source, pages=3)

    first = inspect_document(
        config,
        DocumentInspectionRequest(
            "managed",
            "sample.pdf",
            True,
            content_query="SEARCHABLE pdf",
            content_limit=2,
        ),
    )
    assert first.content_search is not None
    assert first.content_search.status == "COMPLETE"
    assert first.content_search.match_mode == "SUBSTRING_CASEFOLD_NFKC"
    assert first.content_search.matched_item_count == 3
    assert first.content_search.matched_occurrence_count == 3
    assert first.content_search.returned_count == 2
    assert first.content_search.has_more is True
    assert [item.location["page"] for item in first.content_search.matches] == [1, 2]
    assert first.source_sha256 is not None

    second = inspect_document(
        config,
        DocumentInspectionRequest(
            "managed",
            "sample.pdf",
            True,
            content_query="SEARCHABLE pdf",
            content_limit=2,
            content_offset=2,
            expected_source_sha256=first.source_sha256,
        ),
    )
    assert second.content_search is not None
    assert second.content_search.returned_count == 1
    assert second.content_search.has_more is False
    assert [item.location["page"] for item in second.content_search.matches] == [3]
    assert second.document_observation_digest == first.document_observation_digest

    with pytest.raises(DocumentInspectionInputError):
        inspect_document(
            config,
            DocumentInspectionRequest(
                "managed",
                "sample.pdf",
                True,
                content_query="SEARCHABLE pdf",
                content_offset=2,
            ),
        )

    _write_pdf(source, pages=4)
    with pytest.raises(DocumentInspectionSourceChangedError):
        inspect_document(
            config,
            DocumentInspectionRequest(
                "managed",
                "sample.pdf",
                True,
                content_query="SEARCHABLE pdf",
                content_offset=2,
                expected_source_sha256=first.source_sha256,
            ),
        )


def test_cli_json_human_and_noncomplete_failure_contract(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path, _config_value = _config(tmp_path, monkeypatch)
    _write_pdf(tmp_path / "documents" / "sample.pdf")
    (tmp_path / "documents" / "unknown.pdf").write_bytes(b"not a PDF")
    runner = CliRunner()
    command = [
        "--config",
        str(config_path),
        "documents",
        "inspect",
        "--scope",
        "managed",
        "--path",
        "sample.pdf",
        "--yes",
    ]

    encoded = runner.invoke(app, ["--format", "json", *command])
    human = runner.invoke(app, command)
    payload = json.loads(encoded.stdout)
    assert encoded.exit_code == human.exit_code == 0
    assert payload["command"] == "documents.inspect"
    assert payload["status"] == "COMPLETE"
    assert payload["result"]["inspection"]["status"] == "COMPLETE"
    assert "Inspection Status: COMPLETE" in human.stdout
    assert "searchable PDF page" in human.stdout

    located = runner.invoke(
        app,
        [
            "--format",
            "json",
            "--config",
            str(config_path),
            "documents",
            "inspect",
            "--scope",
            "managed",
            "--path",
            "sample.pdf",
            "--yes",
            "--content-query",
            "SEARCHABLE pdf",
        ],
    )
    located_payload = json.loads(located.stdout)
    assert located.exit_code == 0
    assert located_payload["result"]["inspection"]["content_search"]["matched_item_count"] == 1
    assert located_payload["result"]["inspection"]["content_search"]["matches"][0]["location"] == {
        "block": 1,
        "page": 1,
    }

    evidence_command = [
        "--config",
        str(config_path),
        "documents",
        "inspect",
        "--scope",
        "managed",
        "--path",
        "sample.pdf",
        "--yes",
        "--evidence",
        "--content-query",
        "not present",
    ]
    evidence_json = runner.invoke(app, ["--format", "json", *evidence_command])
    evidence_human = runner.invoke(app, evidence_command)
    evidence_payload = json.loads(evidence_json.stdout)
    assert evidence_json.exit_code == evidence_human.exit_code == 0
    selection = evidence_payload["result"]["inspection"]["evidence_selection"]
    assert selection["status"] == "NO_MATCH"
    assert selection["processing_strategy"] == "QUERY_MAP_THEN_BOUNDED_GRAPH_SELECTION"
    assert "Evidence Selection:" in evidence_human.stdout

    unsupported = runner.invoke(
        app,
        [
            "--format",
            "json",
            "--config",
            str(config_path),
            "documents",
            "inspect",
            "--scope",
            "managed",
            "--path",
            "unknown.pdf",
            "--yes",
        ],
    )
    rejected = json.loads(unsupported.stdout)
    assert unsupported.exit_code == 4
    assert rejected["status"] == "UNSUPPORTED_FORMAT"
    assert rejected["result"]["inspection"]["items"] == []

    confirmation = runner.invoke(
        app,
        [
            "--format",
            "json",
            "--config",
            str(config_path),
            "documents",
            "inspect",
            "--scope",
            "managed",
            "--path",
            "sample.pdf",
        ],
    )
    assert confirmation.exit_code == 2
    assert json.loads(confirmation.stdout)["errors"][0]["code"] == (
        "DOCUMENT_INSPECTION_CONFIRMATION_REQUIRED"
    )


def test_product_ingress_rejects_a_source_on_another_device(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root = tmp_path / "documents"
    root.mkdir()
    _write_pdf(root / "sample.pdf")
    bindings = ScopeBindings((ScopeBinding("managed", root),), (str(root),), ("managed",))
    ingress = ProjectOwnedBoundedDocumentIngress(bindings, require_same_device=True)
    adapter = StructuredDocumentParserAdapter(ingress)
    from local_steward.file_agent.runtime import structured_documents

    original = structured_documents.os.fstat
    calls = 0

    def changed_root_device(descriptor: int):
        nonlocal calls
        state = original(descriptor)
        calls += 1
        if calls == 1:
            return SimpleNamespace(st_dev=state.st_dev + 1)
        return state

    monkeypatch.setattr(structured_documents.os, "fstat", changed_root_device)
    observation = adapter.observe({"scope_id": "managed", "relative_path": "sample.pdf"})
    assert observation.status == "UNAVAILABLE"
    assert observation.items == ()


def test_configured_scope_root_symlink_is_rejected(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    target = tmp_path / "target"
    target.mkdir()
    link = tmp_path / "scope-link"
    link.symlink_to(target, target_is_directory=True)
    _write_pdf(target / "sample.pdf")
    _path, config = _config(
        tmp_path,
        monkeypatch,
        [("managed", "managed_root", link, True)],
    )

    with pytest.raises(DocumentInspectionScopeError):
        inspect_document(config, DocumentInspectionRequest("managed", "sample.pdf", True))


def test_product_surface_does_not_expose_backend_or_limit_overrides() -> None:
    help_result = CliRunner().invoke(app, ["documents", "inspect", "--help"])
    assert help_result.exit_code == 0
    assert "--backend" not in help_result.stdout
    assert "--ocr" not in help_result.stdout
    assert "--format-override" not in help_result.stdout
    assert "--source-limit" not in help_result.stdout
