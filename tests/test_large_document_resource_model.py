"""Large-document admission and query-map regressions for NEXT-014."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from types import SimpleNamespace
from typing import Any
import zipfile

import fitz  # type: ignore[import-untyped]
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.chart import BarChart, Reference
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import pytest

import local_steward.file_agent.runtime.structured_documents as structured_documents
from local_steward.file_agent.runtime import (
    DOCUMENT_INGRESS_CHUNK_BYTES,
    MAX_PACKAGE_SOURCE_BYTES,
    MAX_PDF_SOURCE_BYTES,
    ProjectOwnedBoundedDocumentIngress,
    ScopeBinding,
    ScopeBindings,
    StructuredDocumentParserAdapter,
    identify_document_format,
)
from local_steward.file_agent.runtime.structured_documents import _WorkerExecution


def _bindings(tmp_path: Path) -> tuple[Path, ScopeBindings]:
    root = tmp_path / "documents"
    root.mkdir()
    return root, ScopeBindings((ScopeBinding("managed", root),), (str(root),), ("managed",))


def _arguments(path: str, **extra: object) -> dict[str, object]:
    return {"scope_id": "managed", "relative_path": path, **extra}


def _payload() -> dict[str, Any]:
    return {
        "backend_name": "PyMuPDF4LLM",
        "backend_version": "1.28.2",
        "warnings": [],
        "items": [
            {
                "kind": "pdf_page_block",
                "text_or_value": "streamed source",
                "parent": None,
                "location": {"page": 1, "block": 1},
                "extension": {"page": 1},
            }
        ],
    }


@dataclass
class _RecordingWorker:
    staged_paths: list[Path]

    def run(self, source_path: Path) -> _WorkerExecution:
        self.staged_paths.append(source_path)
        assert source_path.stat().st_size > DOCUMENT_INGRESS_CHUNK_BYTES * 2
        return _WorkerExecution("COMPLETE", _payload(), 1, 1024)


@dataclass
class _ProfileWorker:
    backend_name: str
    text: str

    def run(self, _source_path: Path) -> _WorkerExecution:
        kind = "pdf_ocr_text_line" if self.backend_name == "STEWARDPageOCR" else "pdf_page_block"
        payload = {
            "backend_name": self.backend_name,
            "backend_version": "synthetic",
            "warnings": [],
            "items": [
                {
                    "kind": kind,
                    "role": "PARAGRAPH",
                    "text_or_value": self.text,
                    "parent": None,
                    "location": {"page": 1, "line": 1},
                    "extension": {"text_source": "LOCAL_OCR"},
                }
            ],
        }
        return _WorkerExecution("COMPLETE", payload, 1, 1024)


def test_ingress_streams_hashes_and_cleans_one_staged_source(tmp_path: Path) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "large.pdf"
    source.write_bytes(b"%PDF-" + b"x" * (DOCUMENT_INGRESS_CHUNK_BYTES * 2 + 17))
    reads: list[int] = []

    def recording_read(descriptor: int, count: int) -> bytes:
        reads.append(count)
        return structured_documents.os.read(descriptor, count)

    worker = _RecordingWorker([])
    adapter = StructuredDocumentParserAdapter(
        ProjectOwnedBoundedDocumentIngress(bindings, read_bytes=recording_read)
    )
    adapter.worker = worker  # type: ignore[assignment]

    observation = adapter.observe(_arguments("large.pdf"))

    assert observation.status == "COMPLETE"
    assert len(reads) >= 4 and set(reads) == {DOCUMENT_INGRESS_CHUNK_BYTES}
    assert observation.resources.admission_profile == "PAGINATED_STREAM"
    assert observation.resources.source_limit_bytes == MAX_PDF_SOURCE_BYTES
    assert worker.staged_paths and not worker.staged_paths[0].exists()
    assert source.exists()


def test_large_xlsx_evidence_uses_lazy_query_map(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "large.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Facts"
    sheet["A1"] = "DISTINCTIVE_QUERY_MARKER retained lazily"
    sheet["B2"] = "unrelated"
    sheet["C3"] = 42
    workbook.save(source)
    workbook.close()
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    observation = adapter.observe(
        _arguments(
            "large.xlsx",
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="DISTINCTIVE_QUERY_MARKER",
        )
    )

    assert observation.status == "COMPLETE"
    assert observation.backend_name == "STEWARDStreamingMap"
    assert observation.execution is not None
    assert observation.execution.initial_profile == "MAP"
    assert observation.execution.selected_profile == "MAP"
    assert [attempt.profile for attempt in observation.execution.attempts] == ["MAP"]
    assert observation.execution.selection is not None
    assert observation.execution.selection.map_profile == "MAP"
    assert observation.items[0].location == {"sheet": "Facts", "sheet_index": 1, "cell": "A1"}
    assert observation.resources.source_limit_bytes == MAX_PACKAGE_SOURCE_BYTES

    numeric = adapter.observe(
        _arguments(
            "large.xlsx",
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="42",
        )
    )
    assert numeric.status == "COMPLETE"
    assert numeric.items[0].location["cell"] == "C3"
    assert numeric.items[0].text_or_value == "42"


def test_indexed_shared_string_map_streams_only_matching_cells(tmp_path: Path) -> None:
    package = tmp_path / "indexed.xlsx"
    worksheet = (
        b'<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        b"<sheetData><row>"
        b'<c r="A1" t="s"><v>2</v></c>'
        b'<c r="B1" t="s"><v>7</v></c>'
        b'<c r="C1"><v>7</v></c>'
        b"</row></sheetData></worksheet>"
    )
    with zipfile.ZipFile(package, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("xl/worksheets/sheet1.xml", worksheet)
    items: list[dict[str, object]] = []

    with zipfile.ZipFile(package) as archive:
        scanned, stopped = structured_documents._streaming_xlsx_shared_string_cells(
            archive,
            "xl/worksheets/sheet1.xml",
            sheet_name="Facts",
            sheet_index=1,
            shared_matches={7: "distinctive shared text"},
            items=items,
        )

    assert scanned == 1 and stopped is False
    assert items == [
        {
            "kind": "xlsx_cell",
            "role": "TABLE_CELL",
            "text_or_value": "distinctive shared text",
            "parent": "sheet:1",
            "location": {"sheet": "Facts", "sheet_index": 1, "cell": "B1"},
            "extension": {"streaming_query_map": True, "shared_string_index": 7},
        }
    ]


def test_large_xlsx_query_map_includes_comments_and_chart_cache(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "fidelity.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Review"
    sheet.append(["Quarter", "Revenue"])
    sheet.append(["Q1", 42])
    sheet["A1"].comment = Comment("STREAMED_XLSX_COMMENT", "Reviewer")
    chart = BarChart()
    chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=2), titles_from_data=True)
    sheet.add_chart(chart, "D2")
    workbook.save(source)
    workbook.close()
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    comment = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="STREAMED_XLSX_COMMENT",
        )
    )
    chart_result = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="'Review'!B1",
        )
    )

    comment_item = next(item for item in comment.items if item.kind == "xlsx_comment")
    assert comment_item.location["cell"] == "A1"
    assert comment_item.extension is not None
    assert comment_item.extension["author"] == "Reviewer"
    assert any(item.kind == "xlsx_chart" for item in chart_result.items)


def test_large_docx_evidence_streams_xml_paragraphs(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "large.docx"
    content_types = b"""<?xml version="1.0"?>
    <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
      <Override PartName="/word/document.xml"
        ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
    </Types>"""
    document = b"""<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
      <w:body><w:p><w:r><w:t>STREAMED_DOCX_MARKER retained</w:t></w:r></w:p></w:body>
    </w:document>"""
    with zipfile.ZipFile(source, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("word/document.xml", document)
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    observation = adapter.observe(
        _arguments(
            "large.docx",
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="STREAMED_DOCX_MARKER",
        )
    )

    assert observation.status == "COMPLETE"
    assert observation.backend_name == "STEWARDStreamingMap"
    assert observation.items[0].location["block"] == 1
    assert observation.execution is not None
    assert observation.execution.selection is not None
    assert observation.execution.selection.strategy == (
        "STREAMING_QUERY_MAP_THEN_NATIVE_CONTAINER_QUALITY"
    )

    missing = adapter.observe(
        _arguments(
            "large.docx",
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="ABSENT_MARKER",
        )
    )
    assert missing.status == "COMPLETE" and missing.items == ()
    assert missing.execution is not None
    assert len(missing.execution.attempts) == 1


def test_large_docx_query_map_includes_comments_and_footnotes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "auxiliary.docx"
    with zipfile.ZipFile(source, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.'
            'openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>',
        )
        archive.writestr(
            "word/document.xml",
            '<w:document xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>body</w:t>'
            "</w:r></w:p></w:body></w:document>",
        )
        archive.writestr(
            "word/comments.xml",
            '<w:comments xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main"><w:comment w:id="4"><w:p><w:r><w:t>'
            "STREAMED_DOCX_COMMENT</w:t></w:r></w:p></w:comment></w:comments>",
        )
        archive.writestr(
            "word/footnotes.xml",
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main"><w:footnote w:id="2"><w:p><w:r><w:t>'
            "STREAMED_DOCX_FOOTNOTE</w:t></w:r></w:p></w:footnote></w:footnotes>",
        )
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    comment = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="STREAMED_DOCX_COMMENT",
        )
    )
    footnote = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="STREAMED_DOCX_FOOTNOTE",
        )
    )

    assert comment.items[0].kind == "docx_comment"
    assert comment.items[0].location == {"comment": "4"}
    assert footnote.items[0].kind == "docx_footnote"
    assert footnote.items[0].location == {"footnote": "2"}

    broken_source = root / "broken-auxiliary.docx"
    with zipfile.ZipFile(broken_source, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.'
            'openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>',
        )
        archive.writestr("word/document.xml", b"<w:document xmlns:w='x'><w:p/></w:document>")
        archive.writestr("word/comments.xml", b"<w:comments>")
        archive.writestr(
            "word/footnotes.xml",
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/'
            'wordprocessingml/2006/main"><w:footnote w:id="2"><w:p><w:r><w:t>'
            "RECOVERED_FOOTNOTE</w:t></w:r></w:p></w:footnote></w:footnotes>",
        )
    recovered = adapter.observe(
        _arguments(
            broken_source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="RECOVERED_FOOTNOTE",
        )
    )
    assert recovered.status == "COMPLETE"
    assert recovered.items[0].kind == "docx_footnote"
    assert "STREAMING_COMPONENT_MALFORMED:comment" in recovered.warnings


def test_large_pptx_query_map_resolves_notes_and_chart_relationships(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "fidelity.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "STREAMED_PPTX_NOTES"
    chart_data = ChartData()
    chart_data.categories = ["Q1"]
    chart_data.add_series("STREAMED_PPTX_CHART", (42,))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        chart_data,
    )
    presentation.save(source)
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    notes = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="STREAMED_PPTX_NOTES",
        )
    )
    chart_result = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="STREAMED_PPTX_CHART",
        )
    )

    assert notes.items[0].kind == "pptx_speaker_notes"
    assert notes.items[0].location["notes"] == 1
    assert chart_result.items[0].kind == "pptx_chart"
    assert chart_result.items[0].location == {"slide": 1, "chart": 1}


def test_large_epub_map_tolerates_html_chapters_and_accounts_fallback_members(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "book.epub"
    with zipfile.ZipFile(source, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        archive.writestr(
            "META-INF/container.xml",
            '<container xmlns="urn:oasis:names:tc:opendocument:xmlns:container" '
            'version="1.0"><rootfiles><rootfile full-path="OPS/book.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        archive.writestr(
            "OPS/book.opf",
            '<package xmlns="http://www.idpf.org/2007/opf" version="3.0">'
            '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            "<dc:title>Tolerant Book</dc:title></metadata></package>",
        )
        archive.writestr(
            "OPS/strict.xhtml",
            '<html xmlns="http://www.w3.org/1999/xhtml"><body>'
            "<p>ordinary strict chapter</p></body></html>",
        )
        archive.writestr(
            "OPS/legacy.html",
            "<html><body><h2>Legacy</h2><p>TOLERANT_EPUB_MARKER & retained<br>"
            "without strict XHTML</body></html>",
        )
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    observation = adapter.observe(
        _arguments(
            "book.epub",
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="TOLERANT_EPUB_MARKER",
        )
    )

    assert observation.status == "COMPLETE"
    assert observation.backend_name == "STEWARDStreamingMap"
    assert len(observation.items) == 1
    assert observation.items[0].location["section_title"] == "legacy.html"
    assert "EPUB_HTML_TOLERANT_FALLBACK_COUNT:1" in observation.warnings
    assert any(
        warning.startswith("EPUB_HTML_TOLERANT_FALLBACK:1:legacy.html")
        for warning in observation.warnings
    )


def test_large_pdf_query_map_scans_pages_without_building_a_full_graph(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "large.pdf"
    document = fitz.open()
    for page_number in range(1, 4):
        page = document.new_page()
        page.insert_text((72, 72), f"ordinary searchable page text number {page_number}")
    document.save(source)
    document.close()
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    observation = adapter.observe(
        _arguments(
            "large.pdf",
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="ABSENT_MARKER",
        )
    )

    assert observation.status == "COMPLETE" and observation.items == ()
    assert observation.backend_name == "STEWARDStreamingMap"
    assert "STREAMING_QUERY_MAP_PAGES:3" in observation.warnings
    assert "STREAMING_QUERY_MAP_OCR_PAGES:0" in observation.warnings
    assert observation.execution is not None
    assert observation.execution.initial_profile == "MAP"
    assert observation.execution.selected_profile == "MAP"
    assert [attempt.profile for attempt in observation.execution.attempts] == ["MAP"]


def test_page_local_rapidocr_projection_retains_model_authority_and_region(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "scan.pdf"
    document = fitz.open()
    document.new_page()
    document.save(source)
    document.close()

    class FakeRapidOCR:
        def __call__(self, _image_bytes: bytes) -> SimpleNamespace:
            return SimpleNamespace(
                txts=("SCANNED_QUERY_MARKER retained",),
                scores=(0.975,),
                boxes=(((10, 20), (200, 20), (200, 60), (10, 60)),),
            )

    monkeypatch.setattr(structured_documents, "_rapidocr_engine", FakeRapidOCR)

    matched, pages, native_pages, ocr_pages, warnings = (
        structured_documents._pdf_page_text_projection(
            str(source), query="SCANNED_QUERY_MARKER"
        )
    )
    broad, *_counts = structured_documents._pdf_page_text_projection(str(source), query=None)
    streamed = structured_documents._streaming_pdf_query_map(
        str(source), "SCANNED_QUERY_MARKER"
    )

    assert (pages, native_pages, ocr_pages) == (1, 0, 1)
    assert matched[0]["kind"] == "pdf_ocr_page_block"
    assert matched[0]["extension"]["text_source"] == "LOCAL_OCR"
    assert "OCR_ENGINE:RAPIDOCR_ONNXRUNTIME" in warnings
    assert streamed["resource_extension"] == {
        "ocr_backend": "RapidOCR",
        "ocr_version": structured_documents.version("rapidocr"),
        "ocr_page_count": 1,
        "ocr_text_authority": "MODEL_DERIVED",
    }
    assert broad[0]["kind"] == "pdf_ocr_text_line"
    extension = broad[0]["extension"]
    assert extension["ocr_confidence"] == 0.975
    assert extension["normalized_region"] == [
        pytest.approx(10 / 596, abs=0.002),
        pytest.approx(20 / 842, abs=0.002),
        pytest.approx(200 / 596, abs=0.002),
        pytest.approx(60 / 842, abs=0.002),
    ]
    assert broad[0]["node_id"] == "pdf:page:1:ocr-line:1"
    assert extension["visual_region"]["page"] == 1
    assert extension["visual_region"]["coordinate_space"] == "PAGE_POINTS_TOP_LEFT"


def test_scan_heavy_pdf_routes_queries_to_map_and_read_to_page_ocr(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "scan.pdf"
    document = fitz.open()
    document.new_page()
    document.save(source)
    document.close()
    assert source.stat().st_size < structured_documents.STREAMING_QUERY_MAP_THRESHOLD_BYTES
    assert structured_documents._pdf_native_text_probe(source) == (1, 1)

    profiles: list[str] = []
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    def worker_for(
        _self: StructuredDocumentParserAdapter,
        _source_format: str,
        profile: str,
        **_options: object,
    ) -> tuple[str, _ProfileWorker]:
        profiles.append(profile)
        backend = "STEWARDStreamingMap" if profile == "MAP" else "STEWARDPageOCR"
        return backend, _ProfileWorker(backend, "SCANNED_QUERY_MARKER retained")

    monkeypatch.setattr(StructuredDocumentParserAdapter, "_worker_for", worker_for)

    evidence = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="SCANNED_QUERY_MARKER",
        )
    )
    broad = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="READ",
            view="READ",
        )
    )

    assert profiles == ["MAP", "OCR_NATIVE"]
    assert evidence.status == "COMPLETE"
    assert evidence.execution is not None
    assert evidence.execution.selected_profile == "MAP"
    assert [attempt.profile for attempt in evidence.execution.attempts] == ["MAP"]
    assert evidence.execution.selection is not None
    assert evidence.execution.selection.strategy == (
        "STREAMING_QUERY_MAP_PAGE_LOCAL_PROJECTION"
    )
    assert broad.status == "COMPLETE"
    assert broad.backend_name == "STEWARDPageOCR"
    assert broad.execution is not None
    assert broad.execution.selected_profile == "OCR_NATIVE"


def test_small_many_page_pdf_queries_use_map_without_the_byte_threshold(
    tmp_path: Path,
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "many-pages-small.pdf"
    document = fitz.open()
    for page_number in range(1, 271):
        page = document.new_page()
        text = (
            "Sol und Luna transmission marker"
            if page_number == 203
            else f"ordinary searchable page text number {page_number}"
        )
        page.insert_text((72, 72), text)
    document.save(source)
    document.close()
    assert source.stat().st_size < structured_documents.STREAMING_QUERY_MAP_THRESHOLD_BYTES
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    evidence = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="ABSENT_MARKER",
        )
    )
    located = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="LOCATE",
            view="READ",
            content_query="Sol und Luna",
        )
    )

    assert evidence.status == "COMPLETE" and evidence.items == ()
    assert evidence.backend_name == "STEWARDStreamingMap"
    assert evidence.execution is not None
    assert evidence.execution.initial_profile == "MAP"
    assert evidence.execution.selected_profile == "MAP"
    assert located.status == "COMPLETE"
    assert located.backend_name == "STEWARDStreamingMap"
    assert [item.location["page"] for item in located.items] == [203]
    assert located.execution is not None
    assert located.execution.initial_profile == "MAP"
    assert located.execution.selected_profile == "MAP"


def test_large_pdf_query_map_includes_native_annotation_evidence(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root, bindings = _bindings(tmp_path)
    source = root / "annotated-large.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "ordinary searchable page text")
    annotation = page.add_text_annot((100, 100), "reviewer decision: retain provenance")
    annotation.update()
    document.save(source)
    document.close()
    monkeypatch.setattr(structured_documents, "STREAMING_QUERY_MAP_THRESHOLD_BYTES", 1)
    adapter = StructuredDocumentParserAdapter(ProjectOwnedBoundedDocumentIngress(bindings))

    observation = adapter.observe(
        _arguments(
            source.name,
            parser_profile="AUTO",
            intent="EVIDENCE",
            view="READ",
            content_query="retain provenance",
        )
    )

    assert observation.status == "COMPLETE"
    native = next(item for item in observation.items if item.kind == "pdf_annotation")
    assert native.location == {"page": 1, "annotation": 1}
    assert native.text_or_value == "reviewer decision: retain provenance"
    assert "STREAMING_QUERY_MAP_NATIVE_MATCHES:1" in observation.warnings


def test_container_policy_rejects_duplicate_members_and_extreme_compression(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    duplicate = tmp_path / "duplicate.docx"
    with zipfile.ZipFile(duplicate, "w") as archive:
        archive.writestr("word/document.xml", b"<w:document />")
        archive.writestr("[Content_Types].xml", b"<Types />")
        archive.writestr("word/document.xml", b"<w:document />")
    duplicated = identify_document_format(duplicate, duplicate.name)
    assert duplicated.status == "MALFORMED"
    assert duplicated.reason == "DUPLICATE_CONTAINER_MEMBER"

    compressed = tmp_path / "compressed.xlsx"
    monkeypatch.setattr(structured_documents, "MAX_CONTROL_XML_BYTES", 32)
    monkeypatch.setattr(structured_documents, "MAX_PACKAGE_COMPRESSION_RATIO", 10)
    with zipfile.ZipFile(compressed, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", b"<Types />")
        archive.writestr("xl/workbook.xml", b"<workbook />")
        archive.writestr("xl/repeated.xml", b"x" * 4_096)
    ratio_limited = identify_document_format(compressed, compressed.name)
    assert ratio_limited.status == "RESOURCE_LIMIT"
    assert ratio_limited.reason == "ARCHIVE_COMPRESSION_RATIO_LIMIT"


def test_streaming_optional_component_limit_preserves_the_document(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "optional-component-limit.docx"
    with zipfile.ZipFile(source, "w") as archive:
        archive.writestr("word/comments.xml", b"<comments>" + b"x" * 128 + b"</comments>")
    monkeypatch.setattr(structured_documents, "MAX_OPTIONAL_XML_BYTES", 64)
    warnings: list[str] = []

    with zipfile.ZipFile(source) as archive:
        root = structured_documents._streaming_optional_xml_root(
            archive,
            "word/comments.xml",
            component="comments",
            warnings=warnings,
        )

    assert root is None
    assert warnings == ["STREAMING_COMPONENT_RESOURCE_LIMIT:comments"]


def test_format_limits_are_materially_higher_without_heap_sized_reads() -> None:
    assert MAX_PDF_SOURCE_BYTES >= 1024 * 1024 * 1024
    assert MAX_PACKAGE_SOURCE_BYTES >= 512 * 1024 * 1024
    assert DOCUMENT_INGRESS_CHUNK_BYTES <= 1024 * 1024
