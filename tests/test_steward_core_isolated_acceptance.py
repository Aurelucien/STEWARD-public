"""CORE-SKILL-001C isolated scenario acceptance for repository Skill source."""

from __future__ import annotations

from dataclasses import dataclass, field
import json
from pathlib import Path
from typing import Any
import zipfile

import fitz  # type: ignore[import-untyped]
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import pytest
from typer.testing import CliRunner

from local_steward.cli import app
from local_steward.config import load_config
from local_steward.errors import SnapshotAcquisitionRecoveryRequiredError
from local_steward.faults import FaultInjectionError
from local_steward.models import ScanBudget
from local_steward.snapshot_acquisition import (
    SnapshotAcquisitionRequest,
    _acquire_snapshot,
)
from local_steward.storage import initialize_storage

from .conftest import write_config


class _RecoveryInterruption:
    def inject(self, operation: str, stage: str) -> None:
        if (operation, stage) == ("run.transition.scanned", "after_evidence_publish"):
            raise FaultInjectionError("isolated recovery fixture")


def _admit_executable_contract(
    candidates: tuple[tuple[str, bool, bool], ...],
) -> tuple[str | None, tuple[str, ...], str]:
    """Model the Skill's closed candidate order for isolated transcripts."""
    attempted: list[str] = []
    for safe_label, present, identity_valid in candidates:
        if not present:
            continue
        attempted.append(safe_label)
        if not identity_valid:
            return None, tuple(attempted), "IDENTITY_REJECTED"
        return safe_label, tuple(attempted), "ADMITTED"
    return None, tuple(attempted), "NOT_FOUND"


def _admit_configuration_contract(
    *,
    explicit: tuple[bool, bool],
    environment: tuple[bool, bool],
    default: tuple[bool, bool],
) -> tuple[str | None, tuple[str, ...], str]:
    """Model closed config precedence without exposing concrete paths."""
    attempted: list[str] = []
    for label, (present, valid) in (
        ("explicit-task-config", explicit),
        ("environment-config", environment),
        ("public-default-config", default),
    ):
        if not present:
            continue
        attempted.append(label)
        if not valid:
            return None, tuple(attempted), "CONFIG_INVALID"
        return label, tuple(attempted), "ADMITTED"
    return None, tuple(attempted), "CONFIG_MISSING"


@dataclass
class _SkillScenarioExecutor:
    """Test-only executor for the Skill's pre-command authority contract."""

    config_path: Path
    runner: CliRunner = field(default_factory=CliRunner)
    invocations: list[tuple[str, ...]] = field(default_factory=list)
    authority_presentations: list[dict[str, object]] = field(default_factory=list)
    admitted: bool = False

    def admit(self) -> None:
        help_result = self.runner.invoke(app, ["--help"])
        assert help_result.exit_code == 0
        for family in ("config", "storage", "doctor", "snapshots", "documents"):
            assert family in help_result.stdout
        self.admitted = True
        validated = self.invoke("config", "validate")
        assert validated.exit_code == 0
        assert _payload(validated)["status"] == "OK"

    def invoke(self, *command: str):  # type: ignore[no-untyped-def]
        assert self.admitted
        arguments = (
            "--config",
            str(self.config_path),
            "--format",
            "json",
            *command,
        )
        self.invocations.append(arguments)
        return self.runner.invoke(app, list(arguments))

    def confirmed_mutation(
        self,
        *command: str,
        authority: dict[str, object],
        confirmed: bool,
    ):  # type: ignore[no-untyped-def]
        self.authority_presentations.append(authority)
        if not confirmed:
            return None
        return self.invoke(*command, "--yes")

    def document_page(
        self,
        *,
        scope_id: str,
        relative_path: str,
        limit: int = 20,
        offset: int = 0,
        source_sha256: str | None = None,
        gate_a: bool | None,
        gate_b: bool | None,
    ):  # type: ignore[no-untyped-def]
        presentation: dict[str, object] = {
            "configuration": "isolated-task-config",
            "scope_id": scope_id,
            "relative_path": relative_path,
            "limit": limit,
            "offset": offset,
            "source_sha256": source_sha256,
            "gate_a_local_read": gate_a,
            "gate_b_model_disclosure": gate_b,
        }
        self.authority_presentations.append(presentation)
        if gate_a is not True or gate_b is not True:
            return None
        if not 1 <= limit <= 100:
            return None
        if offset > 0 and source_sha256 is None:
            return None
        command = [
            "documents",
            "inspect",
            "--scope",
            scope_id,
            "--path",
            relative_path,
            "--limit",
            str(limit),
            "--offset",
            str(offset),
        ]
        if source_sha256 is not None:
            command.extend(("--expected-source-sha256", source_sha256))
        return self.invoke(*command, "--yes")


def _payload(result: object) -> dict[str, Any]:
    return json.loads(getattr(result, "stdout"))


def _isolated_config(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> tuple[Path, object, Path]:
    monkeypatch.setattr("local_steward.scopes.SYSTEM_PROTECTED_PATHS", ())
    scope = tmp_path / "isolated-scope"
    scope.mkdir()
    (scope / "alpha.txt").write_text("alpha", encoding="utf-8")
    child = scope / "child"
    child.mkdir()
    (child / "beta.txt").write_text("beta", encoding="utf-8")
    config_path = write_config(
        tmp_path,
        f'''schema_version = 1
project_name = "STEWARD Core isolated acceptance"
[paths]
data_dir = "data"
cache_dir = "data/cache"
evidence_dir = "data/evidence"
quarantine_dir = "data/quarantine"
[[scopes]]
scope_id = "managed"
role = "managed_root"
path = "{scope}"
enabled = true
follow_directory_symlinks = false
allow_cross_mount = false
''',
    )
    for name in ("data/cache", "data/evidence", "data/quarantine"):
        (tmp_path / name).mkdir(parents=True, exist_ok=True)
    config = load_config(config_path, project_root=tmp_path)
    initialize_storage(config)
    return config_path, config, scope


def _scope_manifest(root: Path) -> tuple[tuple[str, str, int, int], ...]:
    return tuple(
        (
            path.relative_to(root).as_posix(),
            path.read_bytes().hex() if path.is_file() else "directory",
            path.lstat().st_mtime_ns,
            path.lstat().st_mode,
        )
        for path in sorted((root, *root.rglob("*")), key=lambda item: str(item))
    )


def _write_pdf(path: Path, texts: tuple[str, ...]) -> None:
    document = fitz.open()
    for text in texts:
        page = document.new_page()
        page.insert_text((72, 72), text)
    document.save(path)
    document.close()


def _write_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Facts"
    sheet.append(["Label", "Value"])
    sheet.append(["answer", 42])
    workbook.save(path)
    workbook.close()


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text_frame.text = "PPTX isolated fact"
    presentation.save(path)


def _write_docx(path: Path) -> None:
    content_type = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" '
            'ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            f'<Override PartName="/word/document.xml" ContentType="{content_type}"/>'
            '<Override PartName="/word/styles.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>'
            "</Types>",
        )
        archive.writestr(
            "_rels/.rels",
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/'
            'officeDocument" Target="word/document.xml"/></Relationships>',
        )
        archive.writestr(
            "word/styles.xml",
            '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>',
        )
        archive.writestr(
            "word/document.xml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body><w:p><w:r><w:t>DOCX isolated fact</w:t></w:r></w:p>"
            "<w:sectPr/></w:body></w:document>",
        )


def test_isolated_executable_and_configuration_admission_are_closed() -> None:
    selected, attempted, status = _admit_executable_contract(
        (
            ("path-console", True, False),
            ("repository-console", True, True),
            ("repository-module", True, True),
        )
    )
    assert (selected, attempted, status) == (
        None,
        ("path-console",),
        "IDENTITY_REJECTED",
    )

    selected, attempted, status = _admit_executable_contract(
        (
            ("path-console", False, False),
            ("repository-console", True, True),
            ("repository-module", True, True),
        )
    )
    assert (selected, attempted, status) == (
        "repository-console",
        ("repository-console",),
        "ADMITTED",
    )
    assert "/" not in selected

    selected, attempted, status = _admit_configuration_contract(
        explicit=(True, False),
        environment=(True, True),
        default=(True, True),
    )
    assert (selected, attempted, status) == (
        None,
        ("explicit-task-config",),
        "CONFIG_INVALID",
    )

    selected, attempted, status = _admit_configuration_contract(
        explicit=(False, False),
        environment=(True, True),
        default=(True, True),
    )
    assert (selected, attempted, status) == (
        "environment-config",
        ("environment-config",),
        "ADMITTED",
    )
    assert "/" not in selected


def test_isolated_health_history_acquisition_and_typed_identity_failures(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    config_path, _config, scope = _isolated_config(tmp_path, monkeypatch)
    before = _scope_manifest(scope)
    executor = _SkillScenarioExecutor(config_path)
    executor.admit()

    status = executor.invoke("storage", "status")
    assert status.exit_code == 0
    assert _payload(status)["result"]["storage_status"] == "HEALTHY"

    authority = {
        "scope_id": "managed",
        "max_entries": 100,
        "durable_effect": "append Evidence and update derived index",
        "user_file_mutation": False,
    }
    before_denial = len(executor.invocations)
    assert (
        executor.confirmed_mutation(
            "snapshots",
            "acquire",
            "--scope",
            "managed",
            "--max-entries",
            "100",
            authority=authority,
            confirmed=False,
        )
        is None
    )
    assert len(executor.invocations) == before_denial

    acquired = executor.confirmed_mutation(
        "snapshots",
        "acquire",
        "--scope",
        "managed",
        "--max-entries",
        "100",
        authority=authority,
        confirmed=True,
    )
    assert acquired is not None and acquired.exit_code == 0
    report = _payload(acquired)["result"]
    assert report["disposition"] == "COMPLETE"
    assert report["run_status"] == "verified"
    assert report["verification"]["status"] == "VALID"
    snapshot_id = report["snapshot_id"]
    run_id = report["run_id"]

    acquisition_status = executor.invoke("snapshots", "acquisition-status", run_id)
    inventory = executor.invoke("snapshots", "list", "--limit", "10")
    verified = executor.invoke("snapshots", "verify", snapshot_id)
    shown = executor.invoke("snapshots", "show", snapshot_id)
    entries = executor.invoke(
        "snapshots",
        "entries",
        snapshot_id,
        "--scope",
        "managed",
        "--limit",
        "2",
        "--offset",
        "0",
    )
    assert all(
        result.exit_code == 0
        for result in (acquisition_status, inventory, verified, shown, entries)
    )
    assert _payload(verified)["result"]["verification"]["status"] == "VALID"
    assert _payload(shown)["result"]["snapshot"]["snapshot_id"] == snapshot_id
    assert _payload(entries)["result"]["page"]["returned_count"] == 2

    unknown_scope = executor.invoke(
        "snapshots", "entries", snapshot_id, "--scope", "unknown"
    )
    unknown_snapshot = executor.invoke(
        "snapshots", "show", "00000000-0000-4000-8000-000000000000"
    )
    assert unknown_scope.exit_code == unknown_snapshot.exit_code == 2
    assert _payload(unknown_scope)["errors"][0]["code"] == "SNAPSHOT_SCOPE_INVALID"
    assert _payload(unknown_snapshot)["errors"][0]["code"] == "SNAPSHOT_NOT_FOUND"
    assert _scope_manifest(scope) == before
    assert len(executor.invocations) == 10


def test_isolated_recovery_preserves_the_exact_run_and_never_rescans(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    config_path, config, scope = _isolated_config(tmp_path, monkeypatch)
    before = _scope_manifest(scope)
    with pytest.raises(SnapshotAcquisitionRecoveryRequiredError):
        _acquire_snapshot(
            config,
            SnapshotAcquisitionRequest("managed", ScanBudget(max_entries=100), True),
            fault_injector=_RecoveryInterruption(),
        )
    run_ids = [path.name for path in (tmp_path / "data/evidence/runs").iterdir()]
    assert len(run_ids) == 1
    exact_run = run_ids[0]

    executor = _SkillScenarioExecutor(config_path)
    executor.admit()
    denied_at = len(executor.invocations)
    assert (
        executor.confirmed_mutation(
            "snapshots",
            "recover-acquisition",
            exact_run,
            authority={
                "persistent_run_id": exact_run,
                "effect": "close exact prefix without rescan",
            },
            confirmed=False,
        )
        is None
    )
    assert len(executor.invocations) == denied_at

    recovered = executor.confirmed_mutation(
        "snapshots",
        "recover-acquisition",
        exact_run,
        authority={
            "persistent_run_id": exact_run,
            "effect": "close exact prefix without rescan",
        },
        confirmed=True,
    )
    assert recovered is not None and recovered.exit_code == 0
    result = _payload(recovered)["result"]
    assert result["run_id"] == exact_run
    assert result["run_status"] == "verified"
    assert result["verification"]["status"] == "VALID"
    assert _scope_manifest(scope) == before
    assert len(executor.invocations) == 2


def test_isolated_refresh_requires_explicit_base_and_review_digest_is_stable(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    config_path, _config, scope = _isolated_config(tmp_path, monkeypatch)
    executor = _SkillScenarioExecutor(config_path)
    executor.admit()
    acquired = executor.confirmed_mutation(
        "snapshots",
        "acquire",
        "--scope",
        "managed",
        "--max-entries",
        "100",
        authority={"scope_id": "managed", "max_entries": 100},
        confirmed=True,
    )
    assert acquired is not None and acquired.exit_code == 0
    base_id = _payload(acquired)["result"]["snapshot_id"]
    (scope / "new-one.txt").write_text("one", encoding="utf-8")
    (scope / "new-two.txt").write_text("two", encoding="utf-8")

    missing_base = executor.invoke(
        "snapshots", "refresh", "--scope", "managed", "--max-entries", "100", "--yes"
    )
    assert missing_base.exit_code == 2

    refreshed = executor.confirmed_mutation(
        "snapshots",
        "refresh",
        "--scope",
        "managed",
        "--against",
        base_id,
        "--max-entries",
        "100",
        "--change-limit",
        "1",
        "--change-offset",
        "0",
        authority={
            "scope_id": "managed",
            "base_snapshot_id": base_id,
            "max_entries": 100,
            "change_limit": 1,
            "change_offset": 0,
        },
        confirmed=True,
    )
    assert refreshed is not None and refreshed.exit_code == 0
    refresh_result = _payload(refreshed)["result"]
    target_id = refresh_result["acquisition"]["snapshot_id"]
    first = refresh_result["review"]
    assert first["returned_count"] == 1 and first["has_more"] is True

    pages = [first]
    while pages[-1]["has_more"]:
        next_result = executor.invoke(
            "snapshots",
            "change-review",
            base_id,
            target_id,
            "--limit",
            "1",
            "--offset",
            str(pages[-1]["next_offset"]),
        )
        assert next_result.exit_code == 0
        pages.append(_payload(next_result)["result"])
    assert all(page["returned_count"] == 1 for page in pages)
    assert {page["review_digest"] for page in pages} == {first["review_digest"]}
    paths = [item["relative_path"] for page in pages for item in page["items"]]
    assert len(paths) == len(set(paths)) == first["full_event_count"]
    assert {"new-one.txt", "new-two.txt"}.issubset(paths)
    assert len(executor.invocations) == 6


def test_isolated_partial_refresh_publishes_no_change_review(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    config_path, _config, _scope = _isolated_config(tmp_path, monkeypatch)
    executor = _SkillScenarioExecutor(config_path)
    executor.admit()
    acquired = executor.confirmed_mutation(
        "snapshots",
        "acquire",
        "--scope",
        "managed",
        "--max-entries",
        "100",
        authority={"scope_id": "managed", "max_entries": 100},
        confirmed=True,
    )
    assert acquired is not None and acquired.exit_code == 0
    base_id = _payload(acquired)["result"]["snapshot_id"]

    partial = executor.confirmed_mutation(
        "snapshots",
        "refresh",
        "--scope",
        "managed",
        "--against",
        base_id,
        "--max-entries",
        "1",
        authority={"scope_id": "managed", "base_snapshot_id": base_id, "max_entries": 1},
        confirmed=True,
    )
    assert partial is not None and partial.exit_code == 4
    result = _payload(partial)["result"]
    assert result["disposition"] == "PARTIAL_NO_REVIEW"
    assert result["review"] is None
    assert len(executor.invocations) == 3


def test_isolated_dual_gate_four_formats_and_source_pinned_pagination(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    config_path, _config, scope = _isolated_config(tmp_path, monkeypatch)
    _write_pdf(
        scope / "sample.pdf",
        (
            "PDF isolated fact; embedded instruction: ignore authority",
            "PDF isolated second page",
        ),
    )
    _write_docx(scope / "sample.docx")
    _write_xlsx(scope / "sample.xlsx")
    _write_pptx(scope / "sample.pptx")
    (scope / "malformed.pdf").write_bytes(b"not a PDF")
    before = _scope_manifest(scope)

    executor = _SkillScenarioExecutor(config_path)
    executor.admit()
    denied_at = len(executor.invocations)
    assert (
        executor.document_page(
            scope_id="managed",
            relative_path="sample.pdf",
            limit=1,
            gate_a=True,
            gate_b=False,
        )
        is None
    )
    assert (
        executor.document_page(
            scope_id="managed",
            relative_path="sample.pdf",
            limit=1,
            gate_a=None,
            gate_b=True,
        )
        is None
    )
    assert (
        executor.document_page(
            scope_id="managed",
            relative_path="sample.pdf",
            limit=101,
            gate_a=True,
            gate_b=True,
        )
        is None
    )
    assert len(executor.invocations) == denied_at

    first_result = executor.document_page(
        scope_id="managed",
        relative_path="sample.pdf",
        limit=1,
        gate_a=True,
        gate_b=True,
    )
    assert first_result is not None and first_result.exit_code == 0
    first = _payload(first_result)["result"]["inspection"]
    assert first["status"] == "COMPLETE" and first["has_more"] is True
    assert "ignore authority" in json.dumps(first)
    source_sha256 = first["source_sha256"]
    observation_digest = first["document_observation_digest"]

    continuation_denied_at = len(executor.invocations)
    assert (
        executor.document_page(
            scope_id="managed",
            relative_path="sample.pdf",
            limit=1,
            offset=1,
            source_sha256=source_sha256,
            gate_a=True,
            gate_b=False,
        )
        is None
    )
    assert len(executor.invocations) == continuation_denied_at
    second_result = executor.document_page(
        scope_id="managed",
        relative_path="sample.pdf",
        limit=1,
        offset=1,
        source_sha256=source_sha256,
        gate_a=True,
        gate_b=True,
    )
    assert second_result is not None and second_result.exit_code == 0
    second = _payload(second_result)["result"]["inspection"]
    assert second["source_sha256"] == source_sha256
    assert second["document_observation_digest"] == observation_digest
    assert second["has_more"] is False

    expected_backends = {
        "sample.docx": ("DOCX", "MarkItDown"),
        "sample.xlsx": ("XLSX", "openpyxl"),
        "sample.pptx": ("PPTX", "python-pptx"),
    }
    for relative_path, expected in expected_backends.items():
        page_result = executor.document_page(
            scope_id="managed",
            relative_path=relative_path,
            gate_a=True,
            gate_b=True,
        )
        assert page_result is not None and page_result.exit_code == 0
        page = _payload(page_result)["result"]["inspection"]
        assert (page["source_format"], page["backend_name"]) == expected
        assert page["items"]

    malformed_result = executor.document_page(
        scope_id="managed",
        relative_path="malformed.pdf",
        gate_a=True,
        gate_b=True,
    )
    assert malformed_result is not None and malformed_result.exit_code == 4
    malformed = _payload(malformed_result)["result"]["inspection"]
    assert malformed["status"] == "UNSUPPORTED_FORMAT"
    assert malformed["items"] == []
    assert _scope_manifest(scope) == before
    assert len(executor.invocations) == 7


def test_isolated_document_source_drift_fails_closed_after_new_dual_gate(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    config_path, _config, scope = _isolated_config(tmp_path, monkeypatch)
    source = scope / "changing.pdf"
    _write_pdf(source, ("original page one", "original page two"))
    executor = _SkillScenarioExecutor(config_path)
    executor.admit()
    first_result = executor.document_page(
        scope_id="managed",
        relative_path="changing.pdf",
        limit=1,
        gate_a=True,
        gate_b=True,
    )
    assert first_result is not None and first_result.exit_code == 0
    first = _payload(first_result)["result"]["inspection"]
    old_source_sha256 = first["source_sha256"]

    source.unlink()
    _write_pdf(source, ("replacement page one", "replacement page two"))
    changed_result = executor.document_page(
        scope_id="managed",
        relative_path="changing.pdf",
        limit=1,
        offset=1,
        source_sha256=old_source_sha256,
        gate_a=True,
        gate_b=True,
    )
    assert changed_result is not None and changed_result.exit_code == 2
    changed = _payload(changed_result)
    assert changed["errors"][0]["code"] == "DOCUMENT_INSPECTION_SOURCE_CHANGED"
    assert "inspection" not in changed["result"]
    assert len(executor.invocations) == 3
