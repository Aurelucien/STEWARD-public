"""Acceptance for NEXT-009 region-aware visual document inspection."""

from __future__ import annotations

from hashlib import sha256
from pathlib import Path
import shutil

import fitz  # type: ignore[import-untyped]
from mcp.types import ImageContent
from pptx import Presentation
from pptx.util import Inches
import pytest
from PIL import Image

from local_steward.document_observation import DocumentInspectionRequest, inspect_document
from local_steward.file_agent.runtime import identify_document_format
from local_steward.native_mcp_server import NativeStewardDispatcher, create_codex_host_policy
from local_steward.native_mcp_server.protocol import DOCUMENT_INPUT_SCHEMA, DOCUMENT_TOOL

from .test_document_inspection_product import _config, _write_pdf, _write_pptx
from .test_steward_native_agent_surface import _session


@pytest.fixture(autouse=True)
def _admit_task_owned_temporary_scopes(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr("local_steward.scopes.SYSTEM_PROTECTED_PATHS", ())


def _write_image_only_pdf(path: Path) -> None:
    source = fitz.open()
    page = source.new_page(width=600, height=400)
    page.insert_text((70, 190), "STEWARD VISUAL REGION 2049", fontsize=34)
    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    source.close()
    output = fitz.open()
    target = output.new_page(width=600, height=400)
    target.insert_image(target.rect, stream=pixmap.tobytes("png"))
    output.save(path)
    output.close()


def _write_png(path: Path) -> None:
    document = fitz.open()
    page = document.new_page(width=1_200, height=500)
    page.insert_text((80, 280), "CURRENT IMAGE FACT 7719", fontsize=72)
    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    path.write_bytes(pixmap.tobytes("png"))
    document.close()


@pytest.mark.anyio
async def test_view_returns_ephemeral_mcp_image_without_source_or_directory_mutation(
    tmp_path: Path,
) -> None:
    _config_value, scope, session = _session(tmp_path)
    source = scope / "visual.pdf"
    _write_pdf(source)
    before_source = sha256(source.read_bytes()).hexdigest()
    before_names = sorted(item.name for item in scope.iterdir())

    result = await NativeStewardDispatcher(session, create_codex_host_policy()).dispatch(
        DOCUMENT_TOOL,
        {"action": "VIEW", "absolute_path": str(source), "page": 1, "visual_scale": 1.5},
    )

    assert result.isError is False
    assert len(result.content) == 2
    assert isinstance(result.content[1], ImageContent)
    visual = result.structuredContent["result"]["visual"]
    assert visual["status"] == "COMPLETE"
    assert visual["source_format"] == "PDF"
    assert visual["rendered_page"] == 1
    assert visual["mime_type"] == "image/png"
    assert visual["image_bytes"] > 0
    assert "image_data" not in visual
    assert sha256(source.read_bytes()).hexdigest() == before_source
    assert sorted(item.name for item in scope.iterdir()) == before_names


@pytest.mark.anyio
async def test_view_accepts_a_graph_node_and_returns_its_bounded_region(tmp_path: Path) -> None:
    _config_value, scope, session = _session(tmp_path)
    source = scope / "scan.pdf"
    _write_image_only_pdf(source)
    dispatcher = NativeStewardDispatcher(session, create_codex_host_policy())

    read = await dispatcher.dispatch(
        DOCUMENT_TOOL,
        {"action": "READ", "absolute_path": str(source), "limit": 100},
    )
    items = read.structuredContent["result"]["document"]["items"]
    selected = next(item for item in items if item.get("extension", {}).get("visual_region"))
    view = await dispatcher.dispatch(
        DOCUMENT_TOOL,
        {
            "action": "VIEW",
            "absolute_path": str(source),
            "node_id": selected["node_id"],
        },
    )

    visual = view.structuredContent["result"]["visual"]
    assert visual["status"] == "COMPLETE"
    assert visual["selected_node_id"] == selected["node_id"]
    assert visual["rendered_bbox"] is not None
    assert visual["pixel_width"] > 0
    assert isinstance(view.content[1], ImageContent)


@pytest.mark.anyio
async def test_png_is_discoverable_readable_and_directly_viewable(tmp_path: Path) -> None:
    config, scope, session = _session(tmp_path)
    source = scope / "fact.png"
    _write_png(source)

    identified = identify_document_format(source.read_bytes(), source.name)
    assert identified.source_format == "PNG"
    parsed = inspect_document(
        config,
        DocumentInspectionRequest(
            "managed",
            source.name,
            True,
            parser_profile="AUTO",
            view="READ",
            content_query="CURRENT IMAGE",
        ),
    )
    assert parsed.status == "COMPLETE"
    assert parsed.source_format == "PNG"
    assert parsed.content_search is not None
    assert parsed.content_search.matched_item_count >= 1

    viewed = await NativeStewardDispatcher(session, create_codex_host_policy()).dispatch(
        DOCUMENT_TOOL,
        {"action": "VIEW", "query": "fact.png", "page": 1},
    )
    assert viewed.structuredContent["result"]["visual"]["status"] == "COMPLETE"
    assert viewed.structuredContent["result"]["document_search"]["matched_count"] == 1
    assert isinstance(viewed.content[1], ImageContent)


@pytest.mark.anyio
async def test_large_jpeg_uses_decoder_subsampling_before_rgb_projection(
    tmp_path: Path,
) -> None:
    _config_value, scope, session = _session(tmp_path)
    source = scope / "large.jpg"
    image = Image.new("RGB", (7_500, 4_203), (24, 48, 96))
    image.save(source, "JPEG", quality=90)
    image.close()
    before = sha256(source.read_bytes()).hexdigest()

    result = await NativeStewardDispatcher(session, create_codex_host_policy()).dispatch(
        DOCUMENT_TOOL,
        {"action": "VIEW", "absolute_path": str(source), "page": 1},
    )

    visual = result.structuredContent["result"]["visual"]
    assert visual["status"] == "COMPLETE"
    assert visual["schema_version"] == 2
    assert visual["decode_projection"]["source_pixel_width"] == 7_500
    assert visual["decode_projection"]["source_pixel_height"] == 4_203
    assert visual["decode_projection"]["decoder_output_width"] < 7_500
    assert visual["decode_projection"]["decoder_subsample"] >= 2.0
    assert visual["resource_usage"]["estimated_decode_work_bytes"] > 0
    assert sha256(source.read_bytes()).hexdigest() == before


@pytest.mark.anyio
async def test_pptx_visual_fidelity_uses_ephemeral_libreoffice_projection(
    tmp_path: Path,
) -> None:
    if shutil.which("soffice") is None:
        pytest.skip("LibreOffice headless renderer is unavailable")
    _config_value, scope, session = _session(tmp_path)
    source = scope / "slides.pptx"
    _write_pptx(source)
    before = sha256(source.read_bytes()).hexdigest()

    result = await NativeStewardDispatcher(session, create_codex_host_policy()).dispatch(
        DOCUMENT_TOOL,
        {"action": "VIEW", "absolute_path": str(source), "page": 1},
    )

    visual = result.structuredContent["result"]["visual"]
    assert visual["status"] == "COMPLETE", visual
    assert visual["renderer_name"] == "LibreOffice+PyMuPDF"
    assert isinstance(result.content[1], ImageContent)
    assert sha256(source.read_bytes()).hexdigest() == before


@pytest.mark.anyio
async def test_office_visual_projection_rejects_external_relationships(
    tmp_path: Path,
) -> None:
    _config_value, scope, session = _session(tmp_path)
    source = scope / "external-link.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = "external reference"
    run.hyperlink.address = "https://example.com/untrusted"
    presentation.save(source)
    before = sha256(source.read_bytes()).hexdigest()

    result = await NativeStewardDispatcher(session, create_codex_host_policy()).dispatch(
        DOCUMENT_TOOL,
        {"action": "VIEW", "absolute_path": str(source), "page": 1},
    )

    visual = result.structuredContent["result"]["visual"]
    assert visual["status"] == "UNSUPPORTED_VISUAL"
    assert visual["identification_reason"] == "OFFICE_EXTERNAL_RELATIONSHIPS_PRESENT"
    assert len(result.content) == 1
    assert sha256(source.read_bytes()).hexdigest() == before


def test_formula_view_and_visual_actions_are_additive_to_the_five_tool_surface(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    _path, config = _config(tmp_path, monkeypatch)
    source = tmp_path / "documents" / "formula.pdf"
    _write_pdf(source)
    page = inspect_document(
        config,
        DocumentInspectionRequest(
            "managed", "formula.pdf", True, parser_profile="ENRICHED", view="FORMULAS"
        ),
    )
    assert page.status == "COMPLETE"
    assert page.view == "FORMULAS"
    assert all(item.role == "FORMULA" for item in page.items)
    action_enum = DOCUMENT_INPUT_SCHEMA["properties"]["action"]["enum"]
    assert "VIEW" in action_enum
    assert "EXTRACT_FORMULA" in action_enum


@pytest.mark.anyio
async def test_visual_selectors_are_rejected_outside_view_action(tmp_path: Path) -> None:
    _config_value, scope, session = _session(tmp_path)
    source = scope / "named.pdf"
    _write_pdf(source)
    result = await NativeStewardDispatcher(session, create_codex_host_policy()).dispatch(
        DOCUMENT_TOOL,
        {"action": "READ", "absolute_path": str(source), "page": 1},
    )
    assert result.isError is True
    assert result.structuredContent["error"]["code"] == "STEWARD_NATIVE_ARGUMENT_INVALID"
